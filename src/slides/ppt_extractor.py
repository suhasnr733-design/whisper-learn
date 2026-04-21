"""
Extract text and images from PowerPoint files
"""
from pptx import Presentation
from typing import List, Dict, Optional
from io import BytesIO
import base64

class PPTExtractor:
    def __init__(self):
        self.prs = None
    
    def load_ppt(self, ppt_path: str):
        """Load PowerPoint presentation"""
        self.prs = Presentation(ppt_path)
        return self
    
    def extract_text(self) -> List[Dict]:
        """Extract text from all slides"""
        if self.prs is None:
            raise ValueError("Load PPT first with load_ppt()")
        
        slides_text = []
        
        for slide_idx, slide in enumerate(self.prs.slides):
            slide_content = {
                'slide_num': slide_idx + 1,
                'text': [],
                'notes': None
            }
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        slide_content['text'].append(shape.text)
            
            # Extract speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    slide_content['notes'] = notes_slide.notes_text_frame.text
            
            slide_content['text'] = '\n'.join(slide_content['text'])
            slides_text.append(slide_content)
        
        return slides_text
    
    def extract_images(self, slide_num: Optional[int] = None) -> List[Dict]:
        """Extract images from slides"""
        if self.prs is None:
            raise ValueError("Load PPT first with load_ppt()")
        
        images = []
        slides_to_process = [slide_num - 1] if slide_num else range(len(self.prs.slides))
        
        for slide_idx in slides_to_process:
            slide = self.prs.slides[slide_idx]
            
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image = shape.image
                    img_bytes = image.blob
                    img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                    
                    images.append({
                        'slide_num': slide_idx + 1,
                        'image_index': len(images),
                        'data': img_b64,
                        'ext': image.ext
                    })
        
        return images
    
    def get_slide_count(self) -> int:
        """Get number of slides"""
        if self.prs is None:
            return 0
        return len(self.prs.slides)
    
    def extract_titles(self) -> List[str]:
        """Extract slide titles"""
        if self.prs is None:
            return []
        
        titles = []
        for slide in self.prs.slides:
            if slide.shapes.title:
                titles.append(slide.shapes.title.text)
            else:
                titles.append("")
        
        return titles


class LightweightPPTExtractor:
    """Minimal memory usage PPT extraction"""
    
    def __init__(self, ppt_path: str):
        self.ppt_path = ppt_path
    
    def get_slide_text(self) -> List[str]:
        """Extract text without loading all slides at once"""
        from pptx import Presentation
        
        texts = []
        prs = Presentation(self.ppt_path)
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            texts.append('\n'.join(slide_text))
        
        return texts